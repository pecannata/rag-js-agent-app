#!/usr/bin/env python3
"""
Microsoft PowerPoint Presentation Text Extraction Script
Extracts text from .pptx files using python-pptx.
Includes OCR capability for extracting text from images.
"""

import sys
import json
import argparse
import io
from pathlib import Path

try:
    from pptx import Presentation
except ImportError:
    print(json.dumps({
        "success": False,
        "error": "python-pptx not installed. Please run: pip install python-pptx",
        "text": "",
        "slideCount": 0
    }))
    sys.exit(1)

# OCR dependencies - optional
OCR_AVAILABLE = False
try:
    from PIL import Image
    import pytesseract
    OCR_AVAILABLE = True
except ImportError:
    pass  # OCR will be skipped if dependencies not available

def extract_text_from_images(shape):
    """
    Extract text from images using OCR if available.
    
    Args:
        shape: PowerPoint shape that might contain an image
        
    Returns:
        str: Extracted text from image, empty string if no text or OCR unavailable
    """
    if not OCR_AVAILABLE:
        print(f"OCR not available for shape type: {getattr(shape, 'shape_type', 'unknown')}", file=sys.stderr)
        return ""
    
    try:
        # Import shape types
        from pptx.enum.shapes import MSO_SHAPE_TYPE
        
        print(f"Shape type: {shape.shape_type}, Shape name: {getattr(shape, 'name', 'no name')}", file=sys.stderr)
        
        # Check if shape is a picture
        if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
            print("Found picture shape, attempting OCR...", file=sys.stderr)
            # Get image data
            image_stream = io.BytesIO(shape.image.blob)
            image = Image.open(image_stream)
            
            print(f"Image size: {image.size}, Image mode: {image.mode}", file=sys.stderr)
            
            # Perform OCR with multiple PSM modes for better results
            ocr_configs = [
                '--psm 6',   # Uniform block of text
                '--psm 8',   # Single word
                '--psm 7',   # Single text line
                '--psm 11',  # Sparse text
                '--psm 13'   # Raw line. Treat the image as a single text line
            ]
            
            best_text = ""
            for config in ocr_configs:
                try:
                    ocr_text = pytesseract.image_to_string(image, config=config).strip()
                    if len(ocr_text) > len(best_text):
                        best_text = ocr_text
                except:
                    continue
            
            if best_text:
                print(f"OCR extracted: {best_text[:100]}...", file=sys.stderr)
            else:
                print("No text found in image", file=sys.stderr)
            
            return best_text
        
        # Also check for other shape types that might contain images
        elif hasattr(shape, 'image'):
            print(f"Shape has image attribute, attempting OCR on shape type {shape.shape_type}...", file=sys.stderr)
            try:
                image_stream = io.BytesIO(shape.image.blob)
                image = Image.open(image_stream)
                ocr_text = pytesseract.image_to_string(image, config='--psm 6')
                result = ocr_text.strip()
                if result:
                    print(f"OCR extracted from non-picture shape: {result[:100]}...", file=sys.stderr)
                return result
            except Exception as e:
                print(f"Error processing image in shape: {e}", file=sys.stderr)
                return ""
        
    except Exception as e:
        print(f"OCR processing error: {e}", file=sys.stderr)
        return ""
    
    return ""

def extract_text_from_shape(shape):
    """
    Comprehensive text extraction from a PowerPoint shape.
    
    Args:
        shape: PowerPoint shape object
        
    Returns:
        str: Extracted text from the shape
    """
    text_content = ""
    
    try:
        # Method 1: Direct text attribute
        if hasattr(shape, "text") and shape.text.strip():
            text_content += shape.text.strip() + "\n"
            print(f"    Direct text: {shape.text.strip()[:50]}...", file=sys.stderr)
        
        # Method 2: Text frame (for text boxes and placeholders)
        if hasattr(shape, "text_frame") and shape.text_frame:
            if hasattr(shape.text_frame, "text") and shape.text_frame.text.strip():
                frame_text = shape.text_frame.text.strip()
                if frame_text not in text_content:  # Avoid duplicates
                    text_content += frame_text + "\n"
                    print(f"    Text frame: {frame_text[:50]}...", file=sys.stderr)
            
            # Method 3: Paragraph-level extraction
            if hasattr(shape.text_frame, "paragraphs"):
                for para in shape.text_frame.paragraphs:
                    if hasattr(para, "text") and para.text.strip():
                        para_text = para.text.strip()
                        if para_text not in text_content:  # Avoid duplicates
                            text_content += para_text + "\n"
                            print(f"    Paragraph: {para_text[:50]}...", file=sys.stderr)
                    
                    # Method 4: Run-level extraction (for styled text)
                    if hasattr(para, "runs"):
                        for run in para.runs:
                            if hasattr(run, "text") and run.text.strip():
                                run_text = run.text.strip()
                                if run_text not in text_content:  # Avoid duplicates
                                    text_content += run_text + "\n"
                                    print(f"    Run: {run_text[:50]}...", file=sys.stderr)
        
        # Method 5: Table extraction
        if hasattr(shape, "table"):
            print(f"    Found table with {len(shape.table.rows)} rows", file=sys.stderr)
            for row in shape.table.rows:
                row_text = ""
                for cell in row.cells:
                    if hasattr(cell, "text") and cell.text.strip():
                        row_text += cell.text.strip() + "\t"
                    # Also check cell text frame
                    if hasattr(cell, "text_frame") and cell.text_frame:
                        if hasattr(cell.text_frame, "text") and cell.text_frame.text.strip():
                            cell_frame_text = cell.text_frame.text.strip()
                            if cell_frame_text not in row_text:
                                row_text += cell_frame_text + "\t"
                if row_text.strip():
                    text_content += row_text.strip() + "\n"
                    print(f"    Table row: {row_text.strip()[:50]}...", file=sys.stderr)
        
        # Method 6: Group shape extraction (recursive)
        if hasattr(shape, "shapes"):  # Group shape
            print(f"    Found group with {len(shape.shapes)} sub-shapes", file=sys.stderr)
            for sub_shape in shape.shapes:
                sub_text = extract_text_from_shape(sub_shape)
                if sub_text and sub_text not in text_content:
                    text_content += sub_text
        
        # Method 7: Chart extraction (if text is in chart)
        if hasattr(shape, "chart"):
            print(f"    Found chart shape", file=sys.stderr)
            # Try to extract chart title
            if hasattr(shape.chart, "chart_title") and shape.chart.chart_title:
                if hasattr(shape.chart.chart_title, "text_frame") and shape.chart.chart_title.text_frame:
                    if hasattr(shape.chart.chart_title.text_frame, "text"):
                        chart_title = shape.chart.chart_title.text_frame.text.strip()
                        if chart_title:
                            text_content += chart_title + "\n"
                            print(f"    Chart title: {chart_title[:50]}...", file=sys.stderr)
        
    except Exception as e:
        print(f"    Error extracting text from shape: {e}", file=sys.stderr)
    
    return text_content

def extract_text_from_pptx(pptx_path, slide_by_slide=False):
    """
    Extract text from PowerPoint file including OCR from images.
    
    Args:
        pptx_path (str): Path to the .pptx file
        slide_by_slide (bool): If True, format output for slide-by-slide processing
        
    Returns:
        dict: Result containing success status, text, and metadata
    """
    try:
        presentation = Presentation(pptx_path)
        
        extracted_text = ""
        slide_texts = []
        ocr_text_found = False
        
        print(f"Processing {len(presentation.slides)} slides...", file=sys.stderr)
        print(f"OCR Available: {OCR_AVAILABLE}", file=sys.stderr)
        
        for i, slide in enumerate(presentation.slides):
            print(f"\nProcessing slide {i + 1}...", file=sys.stderr)
            slide_text = ""
            slide_ocr_text = ""
            
            print(f"Slide {i + 1} has {len(slide.shapes)} shapes", file=sys.stderr)
            
            for j, shape in enumerate(slide.shapes):
                print(f"  Shape {j + 1}: Type={getattr(shape, 'shape_type', 'unknown')}", file=sys.stderr)
                
                # Extract comprehensive text from shape
                shape_text = extract_text_from_shape(shape)
                if shape_text.strip():
                    slide_text += shape_text
                
                # Extract text from images using OCR
                if OCR_AVAILABLE:
                    print(f"    Attempting OCR on shape {j + 1}...", file=sys.stderr)
                    ocr_result = extract_text_from_images(shape)
                    if ocr_result:
                        slide_ocr_text += ocr_result + "\n"
                        ocr_text_found = True
                        print(f"    OCR successful: {ocr_result[:50]}...", file=sys.stderr)
                    else:
                        print(f"    No OCR text found", file=sys.stderr)
            
            # Also check slide notes (if any)
            if hasattr(slide, 'notes_slide') and slide.notes_slide:
                print(f"  Found notes slide", file=sys.stderr)
                for notes_shape in slide.notes_slide.shapes:
                    if hasattr(notes_shape, "text") and notes_shape.text.strip():
                        notes_text = notes_shape.text.strip()
                        # Skip the default notes placeholder text
                        if "Click to add notes" not in notes_text:
                            slide_text += "\n[Notes]\n" + notes_text + "\n"
                            print(f"  Notes: {notes_text[:50]}...", file=sys.stderr)
            
            # Combine regular text and OCR text
            combined_text = slide_text
            if slide_ocr_text.strip():
                combined_text += "\n[OCR Text from Images]\n" + slide_ocr_text
            
            combined_text = combined_text.strip()
            if combined_text:  # Only include non-empty text
                slide_texts.append({
                    "slide": i + 1,
                    "text": combined_text,
                    "hasOcrText": bool(slide_ocr_text.strip())
                })
                
                if slide_by_slide:
                    # For slide-by-slide mode, format each slide separately
                    extracted_text += f"\n\n**Slide {i + 1}:**\n{combined_text}\n"
                else:
                    # For standard mode, use existing format
                    extracted_text += f"\n--- Slide {i + 1} ---\n{combined_text}\n"
        
        # Clean up the text
        cleaned_text = extracted_text.strip()
        if cleaned_text:
            # Normalize excessive whitespace but preserve breaks
            lines = []
            for line in cleaned_text.split('\n'):
                cleaned_line = line.strip()
                if cleaned_line or (lines and lines[-1]):  # Keep empty lines only if preceded by content
                    lines.append(cleaned_line)
            cleaned_text = '\n'.join(lines)
        
        return {
            "success": True,
            "text": cleaned_text,
            "slideCount": len(presentation.slides),
            "slideTexts": slide_texts,
            "hasText": bool(cleaned_text.strip()),
            "ocrAvailable": OCR_AVAILABLE,
            "ocrTextFound": ocr_text_found
        }
        
    except FileNotFoundError:
        return {
            "success": False,
            "error": f"PowerPoint file not found: {pptx_path}",
            "text": "",
            "slideCount": 0
        }
    except Exception as e:
        return {
            "success": False,
            "error": f"Unexpected error processing PowerPoint: {str(e)}",
            "text": "",
            "slideCount": 0
        }

def main():
    parser = argparse.ArgumentParser(description='Extract text from Microsoft PowerPoint (.pptx) files')
    parser.add_argument('pptx_path', help='Path to the .pptx file')
    parser.add_argument('--slide-by-slide', action='store_true', help='Format output for slide-by-slide processing')
    
    args = parser.parse_args()
    
    result = extract_text_from_pptx(args.pptx_path, slide_by_slide=args.slide_by_slide)
    print(json.dumps(result, indent=2))

if __name__ == "__main__":
    main()
